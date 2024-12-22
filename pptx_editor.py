import argparse
import os
from pptx import Presentation
from pptx.util import Pt

def parse_arguments():
    """Parse command-line arguments."""
    parser = argparse.ArgumentParser(description="Modify PowerPoint text properties.")
    parser.add_argument("-i", "--input", help="Input file or directory path")
    parser.add_argument("-o", "--output", help="Output file or directory path")
    parser.add_argument("-d", "--directory", action="store_true", help="Indicates the input is a directory")
    parser.add_argument("-s", "--single", action="store_true", help="Indicates the input is a single file")
    parser.add_argument("-f", "--font_size", type=int, help="Font size for all text")
    parser.add_argument("-t", "--title_font_size", type=int, help="Font size for the title text")
    parser.add_argument("-b", "--bold", type=bool, help="Make all text bold (True or False)")

    args = parser.parse_args()

    # Interactive query if no arguments are provided
    if not any(vars(args).values()):
        args = query_user()

    # Validate flags
    if args.directory and args.single:
        parser.error("You cannot use both --directory and --single at the same time.")
    return args

def query_user():
    """Query the user interactively for all options."""
    input_path = input("Enter the input file or directory path: ").strip()
    output_path = input("Enter the output file or directory path: ").strip()
    is_directory = query_yes_no("Is the input a directory? (y/n): ")
    is_single = not is_directory
    font_size = int(input("Enter the font size for all text (or press Enter to skip): ") or 0) or None
    title_font_size = int(input("Enter the font size for title text (or press Enter to skip): ") or 0) or None
    bold = query_yes_no("Should all text be bold? (y/n): ")

    return argparse.Namespace(
        input=input_path,
        output=output_path,
        directory=is_directory,
        single=is_single,
        font_size=font_size,
        title_font_size=title_font_size,
        bold=bold
    )

def query_yes_no(prompt):
    """Prompt the user for a yes/no answer and return a Boolean value."""
    while True:
        response = input(prompt).strip().lower()
        if response in ("y", "yes"):
            return True
        elif response in ("n", "no"):
            return False
        else:
            print("Please answer 'y' or 'n'.")

def process_presentation(file_path, output_path, font_size=None, title_font_size=None, bold=None):
    """Process a single PowerPoint presentation."""
    prs = Presentation(file_path)

    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if font_size:
                            run.font.size = Pt(font_size)
                        if bold is not None:
                            run.font.bold = bold

        # Handle title on the first slide
        if idx == 0 and title_font_size:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text:
                    first_paragraph = shape.text_frame.paragraphs[0]
                    if first_paragraph.runs and first_paragraph.runs[0].font.underline:
                        for run in first_paragraph.runs:
                            run.font.size = Pt(title_font_size)
                        break

    # Save the updated presentation
    prs.save(output_path)
    print("Success: {} processed and saved to {}".format(os.path.basename(file_path), output_path))

def process_directory(input_dir, output_dir, font_size=None, title_font_size=None, bold=None):
    """Process all PowerPoint files in a directory."""
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    for file_name in os.listdir(input_dir):
        if file_name.endswith(".pptx"):
            input_path = os.path.join(input_dir, file_name)
            output_path = os.path.join(output_dir, file_name)
            process_presentation(input_path, output_path, font_size, title_font_size, bold)
        else: 
            print("Skipping non-PPTX files: {}".format(file_name))

def main():
    args = parse_arguments()

    # Process based on mode (single or directory)
    if args.single:
        process_presentation(args.input, args.output, args.font_size, args.title_font_size, args.bold)
    elif args.directory:
        process_directory(args.input, args.output, args.font_size, args.title_font_size, args.bold)

if __name__ == "__main__":
    try:
        main()
    except KeyboardInterrupt:
        print("\nProgram interrupted. Exiting...")
        sys.exit(0)
